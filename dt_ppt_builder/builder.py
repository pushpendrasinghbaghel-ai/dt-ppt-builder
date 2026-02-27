"""
Orchestrator: loads config + data → calls slide builders → saves PPTX.

Usage (CLI):
    from dt_ppt_builder.builder import build
    build("configs/example/config.yaml")

Usage (programmatic / MCP):
    from dt_ppt_builder.builder import build_from_dict
    pptx_bytes = build_from_dict(cfg_dict, requirements_list)
"""
import io
import os
import json
import yaml

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

from . import slide_builder as sb


# ─────────────────────────────────────────────────────────────────────────────
# Template loading — strips all existing slides cleanly (no corruption)
# ─────────────────────────────────────────────────────────────────────────────
def _load_template_clean(path: str) -> Presentation:
    # .potx files need conversion: rename to .pptx so python-pptx accepts them
    import tempfile, shutil
    ext = os.path.splitext(path)[1].lower()
    if ext == ".potx":
        tmp = os.path.join(tempfile.gettempdir(), "dt_ppt_tmp.pptx")
        shutil.copy2(path, tmp)
        # Patch content type inside the ZIP from template → presentation
        import zipfile, io
        buf = io.BytesIO()
        with zipfile.ZipFile(tmp, "r") as zin, zipfile.ZipFile(buf, "w") as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                        b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                    )
                zout.writestr(item, data)
        with open(tmp, "wb") as f:
            f.write(buf.getvalue())
        path = tmp
    prs = Presentation(path)
    xml_slides    = prs.slides._sldIdLst
    slide_id_list = list(xml_slides)
    for sId in slide_id_list:
        rId = sId.get(qn("r:id"))
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        xml_slides.remove(sId)
    assert len(prs.slides) == 0, "Slides not cleared!"
    print(f"  Template loaded: {os.path.basename(path)}")
    return prs


# ─────────────────────────────────────────────────────────────────────────────
# Layout map — maps friendly keys → slide layout objects
# ─────────────────────────────────────────────────────────────────────────────
def _layout_map(prs, cfg) -> dict:
    """
    cfg['layout_indices'] overrides defaults. Keys:
        title_center, title_content, two_img
    """
    defaults = {
        "title_center":  11,   # Title+eyebrow-only_centered  (Perform 2026)
        "title_content":  2,   # Title+eyebrow+content_center
        "two_img":       19,   # 2 images + captions
    }
    overrides = cfg.get("layout_indices", {})
    resolved  = {**defaults, **overrides}
    layouts   = prs.slide_layouts
    SL = {}
    for key, idx in resolved.items():
        if idx < len(layouts):
            SL[key] = layouts[idx]
        else:
            # Graceful fallback: use layout 0
            print(f"  [WARN] layout_indices.{key}={idx} out of range "
                  f"(template has {len(layouts)} layouts); using 0")
            SL[key] = layouts[0]
    return SL


# ─────────────────────────────────────────────────────────────────────────────
# Domain summary helper
# ─────────────────────────────────────────────────────────────────────────────
def _domain_summary(reqs):
    """Return (name, total, now, partial, roadmap) from a list of req dicts."""
    now     = sum(1 for r in reqs if "✅" in r.get("status","") or "Now"     in r.get("status",""))
    partial = sum(1 for r in reqs if "⚡" in r.get("status","") or "Partial" in r.get("status",""))
    roadmap = sum(1 for r in reqs if "🗺" in r.get("status","") or "Roadmap" in r.get("status",""))
    return {"total": len(reqs), "now": now, "partial": partial, "roadmap": roadmap}


# ─────────────────────────────────────────────────────────────────────────────
# Main build function
# ─────────────────────────────────────────────────────────────────────────────
def build(config_path: str):
    """Generate a full branded PPTX from a YAML config + JSON requirements."""
    config_dir = os.path.dirname(os.path.abspath(config_path))

    # ── Load config ──────────────────────────────────────────────────────────
    with open(config_path, encoding="utf-8") as f:
        cfg = yaml.safe_load(f)

    # ── Load requirements data ───────────────────────────────────────────────
    req_path = cfg.get("requirements_file") or os.path.join(config_dir, "requirements.json")
    if not os.path.isabs(req_path):
        req_path = os.path.join(config_dir, req_path)

    with open(req_path, encoding="utf-8") as f:
        req_data = json.load(f)   # list of domain dicts: {name, description, reqs:[...]}

    # ── Load template ────────────────────────────────────────────────────────
    template_path = cfg["template"]
    prs = _load_template_clean(template_path)
    SL  = _layout_map(prs, cfg)

    # ── Slide sequence ───────────────────────────────────────────────────────
    print("  Building slides…")

    # 1. Title slide
    sb.title_slide(prs, SL, cfg)
    print("    ✓ Title slide")

    # 2. Coverage matrix
    domains_summary = []
    for d in req_data:
        summary = _domain_summary(d["reqs"])
        summary["name"] = d["name"]
        domains_summary.append(summary)
    sb.coverage_slide(prs, SL, cfg, domains_summary)
    print("    ✓ Coverage matrix slide")

    # 3. Instrumentation / landing screenshot (if configured)
    landing_bullets = cfg.get("landing_bullets")
    if landing_bullets or cfg.get("images", {}).get("landing"):
        sb.instrumentation_slide(prs, SL, cfg,
                                  img_key="landing",
                                  title=cfg.get("landing_title",
                                                "AI Observability — Application View"),
                                  bullets=landing_bullets or [])
        print("    ✓ Instrumentation slide")

    # 4. One slide per domain
    for d in req_data:
        sb.domain_slide(prs, SL, cfg,
                         domain_label=d["name"],
                         reqs=[(r["requirement"], r["description"],
                                r.get("status", ""), r.get("signal", ""))
                               for r in d["reqs"]],
                         description=d.get("description", ""))
        print(f"    ✓ Domain slide: {d['name']}")

    # 5. Screenshot slides (defined in config as a list)
    for ss in cfg.get("screenshot_slides", []):
        if ss.get("type") == "two_image":
            sb.two_image_slide(prs, SL, cfg,
                                title=ss["title"],
                                left_key=ss["left_key"],   left_caption=ss["left_caption"],
                                right_key=ss["right_key"],  right_caption=ss["right_caption"])
            print(f"    ✓ Two-image slide: {ss['title']}")
        elif ss.get("type") == "single":
            sb.instrumentation_slide(prs, SL, cfg,
                                      img_key=ss["img_key"],
                                      title=ss["title"],
                                      bullets=ss.get("bullets", []))
            print(f"    ✓ Single-image slide: {ss['title']}")

    # 6. GCC / regulatory slide (optional)
    gcc_cfg = cfg.get("gcc_slide")
    if gcc_cfg:
        gcc_reqs_raw = gcc_cfg.get("reqs", [])
        gcc_reqs = [(r["requirement"], r["description"],
                     r.get("status", ""), r.get("signal", ""))
                    for r in gcc_reqs_raw]
        sb.gcc_slide(prs, SL, cfg, gcc_reqs,
                     title=gcc_cfg.get("title"),
                     eyebrow=gcc_cfg.get("eyebrow"))
        print("    ✓ GCC slide")

    # 7. Closing slide
    sb.closing_slide(prs, SL, cfg, message=cfg.get("closing_message"))
    print("    ✓ Closing slide")

    # ── Save ─────────────────────────────────────────────────────────────────
    output = cfg["output"]
    os.makedirs(os.path.dirname(output), exist_ok=True)
    prs.save(output)
    size_mb = os.path.getsize(output) / 1_048_576
    print(f"\n  ✅ Saved: {output}  ({size_mb:.1f} MB, {len(prs.slides)} slides)")
    return output


# ─────────────────────────────────────────────────────────────────────────────
# Programmatic interface (config dict + requirements list → bytes)
# ─────────────────────────────────────────────────────────────────────────────
def _build_prs(cfg: dict, req_data: list) -> Presentation:
    """
    Core builder: takes a config dict and a requirements list,
    returns a python-pptx Presentation object.
    """
    template_path = cfg["template"]
    prs = _load_template_clean(template_path)
    SL  = _layout_map(prs, cfg)

    # 1. Title slide
    sb.title_slide(prs, SL, cfg)

    # 2. Coverage matrix
    domains_summary = []
    for d in req_data:
        summary = _domain_summary(d["reqs"])
        summary["name"] = d["name"]
        domains_summary.append(summary)
    sb.coverage_slide(prs, SL, cfg, domains_summary)

    # 3. Instrumentation / landing slide
    landing_bullets = cfg.get("landing_bullets")
    if landing_bullets or cfg.get("images", {}).get("landing"):
        sb.instrumentation_slide(prs, SL, cfg,
                                  img_key="landing",
                                  title=cfg.get("landing_title",
                                                "AI Observability — Application View"),
                                  bullets=landing_bullets or [])

    # 4. One slide per domain
    for d in req_data:
        sb.domain_slide(prs, SL, cfg,
                         domain_label=d["name"],
                         reqs=[(r["requirement"], r["description"],
                                r.get("status", ""), r.get("signal", ""))
                               for r in d["reqs"]],
                         description=d.get("description", ""))

    # 5. Screenshot slides
    for ss in cfg.get("screenshot_slides", []):
        if ss.get("type") == "two_image":
            sb.two_image_slide(prs, SL, cfg,
                                title=ss["title"],
                                left_key=ss["left_key"],   left_caption=ss["left_caption"],
                                right_key=ss["right_key"],  right_caption=ss["right_caption"])
        elif ss.get("type") == "single":
            sb.instrumentation_slide(prs, SL, cfg,
                                      img_key=ss["img_key"],
                                      title=ss["title"],
                                      bullets=ss.get("bullets", []))

    # 6. GCC / regulatory slide
    gcc_cfg = cfg.get("gcc_slide")
    if gcc_cfg:
        gcc_reqs_raw = gcc_cfg.get("reqs", [])
        gcc_reqs = [(r["requirement"], r["description"],
                     r.get("status", ""), r.get("signal", ""))
                    for r in gcc_reqs_raw]
        sb.gcc_slide(prs, SL, cfg, gcc_reqs,
                     title=gcc_cfg.get("title"),
                     eyebrow=gcc_cfg.get("eyebrow"))

    # 7. Closing slide
    sb.closing_slide(prs, SL, cfg, message=cfg.get("closing_message"))

    return prs


def build_from_dict(cfg: dict, req_data: list) -> bytes:
    """
    Build a PPTX from a config dict + requirements list.
    Returns the .pptx file content as bytes.
    """
    prs = _build_prs(cfg, req_data)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_and_save(cfg: dict, req_data: list, output_path: str) -> str:
    """
    Build a PPTX and save to disk.  Returns the output path.
    """
    prs = _build_prs(cfg, req_data)
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    size_mb = os.path.getsize(output_path) / 1_048_576
    return f"{output_path} ({size_mb:.1f} MB, {len(prs.slides)} slides)"


# ─────────────────────────────────────────────────────────────────────────────
# Generic / dynamic builder (content-driven — no pre-built config required)
# ─────────────────────────────────────────────────────────────────────────────
def build_generic(template_path: str, slides: list[dict],
                  output_path: str, layout_indices: dict | None = None) -> str:
    """
    Build a branded PPTX from an array of slide spec dicts.

    Each slide dict must have a "type" key (title, section, bullets, table,
    two_column, text, image, comparison, closing) plus type-specific fields.

    Args:
        template_path:  Path to .pptx or .potx template
        slides:         List of slide spec dicts
        output_path:    Where to save the output .pptx
        layout_indices: Optional dict overriding layout index mapping

    Returns: summary string "path (size, N slides)"
    """
    from . import generic_slides as gs

    prs = _load_template_clean(template_path)
    cfg = {"layout_indices": layout_indices or {}}
    SL  = _layout_map(prs, cfg)

    gs.render_all(prs, SL, slides)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    size_mb = os.path.getsize(output_path) / 1_048_576
    return f"{output_path} ({size_mb:.1f} MB, {len(prs.slides)} slides)"


def build_generic_bytes(template_path: str, slides: list[dict],
                        layout_indices: dict | None = None) -> bytes:
    """
    Build a branded PPTX from slide specs, return bytes (no disk write).
    """
    from . import generic_slides as gs

    prs = _load_template_clean(template_path)
    cfg = {"layout_indices": layout_indices or {}}
    SL  = _layout_map(prs, cfg)

    gs.render_all(prs, SL, slides)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
