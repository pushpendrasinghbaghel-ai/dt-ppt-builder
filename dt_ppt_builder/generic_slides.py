"""
Generic slide renderer — content-driven slide generation.

Each function takes a Presentation + layout dict + a slide spec dict
and renders one slide. The slide spec is a JSON-friendly dict that
Copilot (or any caller) provides.

Supported slide types:
  title       —  Title + subtitle + optional contact
  section     —  Bold section-divider / chapter header
  bullets     —  Title + bullet points
  table       —  Title + arbitrary table (headers + rows)
  two_column  —  Title + two columns of bullets
  text        —  Title + free-form body text
  image       —  Title + image + optional caption
  comparison  —  Title + two side-by-side comparison blocks
  closing     —  Closing message + contact
  hero        —  Branded hero slide (brand, headline, subtext, tagline)
  card_grid   —  Grid of icon+title+description cards (2x2 or 3x2)
  icon_bullets — Eyebrow + subtitle + icon-prefixed bullets + optional right image
  split_panel —  Left icon-bullets + right panel with grouped content
  two_image   —  Two screenshots side-by-side with captions
  value_props —  Stacked value proposition items (icon + bold title + description)
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

from .brand import WHITE, TEAL, GREEN, ORANGE, GRAY, DGRAY, DDGRAY, DTDARK, RGBColor
from .helpers import set_ph, txb, para_block


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────
def _new(prs, SL, key):
    return prs.slides.add_slide(SL[key])


# Map user-friendly color names → RGBColor
_NAMED_COLORS = {
    "white": WHITE, "teal": TEAL, "green": GREEN,
    "orange": ORANGE, "gray": GRAY,
}

def _resolve_color(val, default=WHITE):
    if val is None:
        return default
    if isinstance(val, str):
        return _NAMED_COLORS.get(val.lower(), default)
    return default


# ─────────────────────────────────────────────────────────────────────────────
# 1. Title slide
# ─────────────────────────────────────────────────────────────────────────────
def render_title(prs, SL, spec):
    sl = _new(prs, SL, "title_center")
    set_ph(sl, 0, spec.get("title", ""),
           size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    set_ph(sl, 1, spec.get("subtitle", ""),
           size=20, color=TEAL, align=PP_ALIGN.CENTER)
    contact = spec.get("contact", "")
    if contact:
        txb(sl, contact, 3.5, 5.6, 7.0, 0.5,
            size=11, color=TEAL, align=PP_ALIGN.CENTER)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 2. Section / chapter divider
# ─────────────────────────────────────────────────────────────────────────────
def render_section(prs, SL, spec):
    sl = _new(prs, SL, "title_center")
    set_ph(sl, 0, spec.get("title", ""),
           size=30, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    sub = spec.get("subtitle", "")
    if sub:
        set_ph(sl, 1, sub, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 3. Bullet-point slide
# ─────────────────────────────────────────────────────────────────────────────
def render_bullets(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)
    eyebrow = spec.get("eyebrow", "")
    if eyebrow:
        set_ph(sl, 1, eyebrow, size=10, color=TEAL, italic=True)
    bullets = spec.get("bullets", [])
    para_block(sl, bullets, 0.7, 2.0, 11.5, 5.2, size=12, color=WHITE)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 4. Table slide
# ─────────────────────────────────────────────────────────────────────────────
def render_table(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)
    eyebrow = spec.get("eyebrow", "")
    if eyebrow:
        set_ph(sl, 1, eyebrow, size=10, color=TEAL, italic=True)

    columns = spec.get("columns", [])
    rows    = spec.get("rows", [])
    if not columns or not rows:
        return sl

    n_cols = len(columns)
    n_rows = len(rows)
    col_w = 12.0 / n_cols

    tbl_shape = sl.shapes.add_table(
        n_rows + 1, n_cols,
        Inches(0.7), Inches(2.2), Inches(12.0), Inches(5.0))
    tbl = tbl_shape.table

    for c in range(n_cols):
        tbl.columns[c].width = Inches(col_w)

    # Header
    for c, h_txt in enumerate(columns):
        cell = tbl.cell(0, c)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(h_txt)
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = TEAL
        cell.fill.solid(); cell.fill.fore_color.rgb = DTDARK

    # Data
    for ri, row in enumerate(rows):
        bg = DGRAY if ri % 2 == 0 else DDGRAY
        for c in range(min(n_cols, len(row))):
            cell = tbl.cell(ri + 1, c)
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(row[c])
            r.font.size = Pt(9); r.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 5. Two-column slide
# ─────────────────────────────────────────────────────────────────────────────
def render_two_column(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)
    eyebrow = spec.get("eyebrow", "")
    if eyebrow:
        set_ph(sl, 1, eyebrow, size=10, color=TEAL, italic=True)

    # Left column
    left_hdr     = spec.get("left_header", "")
    left_bullets = spec.get("left_bullets", [])
    para_block(sl, left_bullets, 0.5, 2.2, 5.8, 4.8,
               size=11, color=WHITE, hdr=left_hdr, hdr_color=TEAL, hdr_size=13)

    # Right column
    right_hdr     = spec.get("right_header", "")
    right_bullets = spec.get("right_bullets", [])
    para_block(sl, right_bullets, 6.8, 2.2, 5.8, 4.8,
               size=11, color=WHITE, hdr=right_hdr, hdr_color=TEAL, hdr_size=13)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 6. Free-text slide
# ─────────────────────────────────────────────────────────────────────────────
def render_text(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)
    eyebrow = spec.get("eyebrow", "")
    if eyebrow:
        set_ph(sl, 1, eyebrow, size=10, color=TEAL, italic=True)
    body = spec.get("body", "")
    txb(sl, body, 0.7, 2.0, 11.5, 5.2, size=12, color=WHITE)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 7. Image slide
# ─────────────────────────────────────────────────────────────────────────────
def render_image(prs, SL, spec):
    import os
    sl = _new(prs, SL, "title_content")
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)
    img_path = spec.get("image_path", "")
    caption  = spec.get("caption", "")
    if img_path and os.path.exists(img_path):
        sl.shapes.add_picture(img_path,
                              Inches(1.5), Inches(1.8), Inches(10.0), Inches(5.0))
    if caption:
        txb(sl, caption, 1.5, 6.9, 10.0, 0.4,
            size=9, color=GRAY, align=PP_ALIGN.CENTER)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 8. Comparison slide (side-by-side)
# ─────────────────────────────────────────────────────────────────────────────
def render_comparison(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)

    items = spec.get("items", [])  # list of {label, bullets}
    n = len(items) or 1
    col_w = 12.0 / n
    for i, item in enumerate(items):
        x = 0.5 + i * col_w
        para_block(sl, item.get("bullets", []),
                   x, 2.2, col_w - 0.3, 4.8,
                   size=11, color=WHITE,
                   hdr=item.get("label", ""), hdr_color=TEAL, hdr_size=13)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 9. Closing slide
# ─────────────────────────────────────────────────────────────────────────────
def render_closing(prs, SL, spec):
    sl = _new(prs, SL, "title_center")
    set_ph(sl, 0, spec.get("message", "Thank you"),
           size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    contact = spec.get("contact", "")
    if contact:
        txb(sl, contact, 3.5, 5.6, 7.0, 0.5,
            size=11, color=TEAL, align=PP_ALIGN.CENTER)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 10. Hero slide — branded opener (brand word, headline, sub-text, tagline)
#     spec: {brand?, headline, sub_headline?, tagline?, footer?}
# ─────────────────────────────────────────────────────────────────────────────
def render_hero(prs, SL, spec):
    sl = _new(prs, SL, "title_center")
    # Brand word
    brand = spec.get("brand", "dynatrace")
    txb(sl, brand, 0.6, 0.4, 3.0, 0.5, size=16, bold=False, color=WHITE)
    # Main headline
    txb(sl, spec.get("headline", ""), 0.6, 1.8, 6.0, 1.2,
        size=44, bold=True, color=WHITE)
    # Sub-headline
    sub = spec.get("sub_headline", "")
    if sub:
        txb(sl, sub, 0.6, 2.85, 6.0, 0.8,
            size=28, bold=False, color=WHITE)
    # Tagline
    tag = spec.get("tagline", "")
    if tag:
        txb(sl, tag, 0.6, 3.8, 6.0, 0.8,
            size=14, color=GRAY)
    # Footer
    footer = spec.get("footer", "")
    if footer:
        txb(sl, footer, 0.5, 5.25, 9.0, 0.3, size=8, color=GRAY)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 11. Card grid — 2x2 or 3x2 grid of cards with colored top-bar
#     spec: {eyebrow, title, cards: [{icon?, title, description}], footer?}
# ─────────────────────────────────────────────────────────────────────────────
_CARD_COLORS = [TEAL, GREEN, ORANGE, RGBColor(0x9B, 0x59, 0xB6),
                RGBColor(0xE7, 0x4C, 0x3C), RGBColor(0x34, 0x95, 0xDB)]

def _card(sl, x, y, w, h, icon, title, desc, bar_color):
    """Draw a single card with color bar, icon, title, description."""
    # Background rectangle
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1A, 0x24, 0x40)
    bg.line.fill.background()
    # Color bar at top
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    # Icon (emoji)
    if icon:
        txb(sl, icon, x + 0.2, y + 0.2, 0.4, 0.4, size=18, color=WHITE)
    # Title
    txb(sl, title, x + 0.65 if icon else x + 0.2, y + 0.2, w - 0.85, 0.4,
        size=12, bold=True, color=WHITE)
    # Description
    txb(sl, desc, x + 0.2, y + 0.65, w - 0.4, h - 0.8,
        size=9, color=GRAY)


def render_card_grid(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)
    eyebrow = spec.get("eyebrow", "")
    if eyebrow:
        set_ph(sl, 1, eyebrow, size=10, color=TEAL, italic=True)

    cards = spec.get("cards", [])
    n = len(cards)
    cols = 3 if n > 4 else 2
    rows = (n + cols - 1) // cols

    card_w = 12.0 / cols - 0.2
    card_h = 4.5 / rows - 0.15
    start_y = 2.0

    for i, c in enumerate(cards):
        col = i % cols
        row = i // cols
        x = 0.6 + col * (card_w + 0.2)
        y = start_y + row * (card_h + 0.15)
        _card(sl, x, y, card_w, card_h,
              c.get("icon", ""), c.get("title", ""), c.get("description", ""),
              _CARD_COLORS[i % len(_CARD_COLORS)])

    footer = spec.get("footer", "")
    if footer:
        txb(sl, footer, 0.5, 5.25, 9.0, 0.3, size=8, color=GRAY)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 12. Icon bullets — eyebrow + subtitle + checkmark-prefixed bullets + optional image
#     spec: {eyebrow, title, subtitle?, bullets: [str], image_path?, image_caption?}
# ─────────────────────────────────────────────────────────────────────────────
def render_icon_bullets(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    eyebrow = spec.get("eyebrow", "")
    if eyebrow:
        set_ph(sl, 1, eyebrow, size=10, color=TEAL, italic=True)
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)

    subtitle = spec.get("subtitle", "")
    if subtitle:
        txb(sl, subtitle, 0.6, 1.5, 5.5, 0.7, size=10, color=GRAY)

    bullets = spec.get("bullets", [])
    img_path = spec.get("image_path", "")
    has_img = img_path and os.path.isfile(img_path)

    bullet_w = 5.0 if has_img else 11.5
    y_start = 2.0 if not subtitle else 2.2
    spacing = 0.42

    for i, b in enumerate(bullets):
        y = y_start + i * spacing
        # Checkmark icon
        txb(sl, "✓", 0.6, y, 0.3, 0.35, size=11, bold=True, color=TEAL)
        # Bullet text
        txb(sl, b, 0.95, y, bullet_w, 0.35, size=10, color=WHITE)

    if has_img:
        sl.shapes.add_picture(img_path,
                              Inches(5.8), Inches(1.0),
                              Inches(3.9), Inches(3.6))
        cap = spec.get("image_caption", "")
        if cap:
            txb(sl, cap, 5.8, 4.7, 3.9, 0.25,
                size=8, color=GRAY, align=PP_ALIGN.CENTER)

    footer = spec.get("footer", "")
    if footer:
        txb(sl, footer, 0.5, 5.25, 9.0, 0.3, size=8, color=GRAY)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 13. Split panel — left icon-bullets + right grouped panel
#     spec: {eyebrow, title, subtitle?, bullets: [str],
#            panel_title?, panel_items: [{text}], footer?}
# ─────────────────────────────────────────────────────────────────────────────
def render_split_panel(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    eyebrow = spec.get("eyebrow", "")
    if eyebrow:
        set_ph(sl, 1, eyebrow, size=10, color=TEAL, italic=True)
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)

    subtitle = spec.get("subtitle", "")
    if subtitle:
        txb(sl, subtitle, 0.6, 0.8, 5.0, 0.7, size=13, color=WHITE)

    # Left: icon bullets
    bullets = spec.get("bullets", [])
    y_start = 1.8
    spacing = 0.42
    for i, b in enumerate(bullets):
        y = y_start + i * spacing
        txb(sl, "✓", 0.6, y, 0.3, 0.35, size=11, bold=True, color=TEAL)
        txb(sl, b, 0.95, y, 4.6, 0.35, size=10, color=WHITE)

    # Right: panel
    panel_items = spec.get("panel_items", [])
    panel_title = spec.get("panel_title", "")
    px, py, pw = 5.8, 1.8, 3.9
    ph = len(panel_items) * 0.42 + 0.6

    # Panel background
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(px), Inches(py), Inches(pw), Inches(ph))
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1A, 0x24, 0x40)
    bg.line.fill.background()
    # Panel color bar
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(px), Inches(py), Inches(pw), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    if panel_title:
        txb(sl, panel_title, px + 0.2, py + 0.15, pw - 0.4, 0.35,
            size=12, bold=True, color=WHITE)
    for i, item in enumerate(panel_items):
        text = item if isinstance(item, str) else item.get("text", "")
        iy = py + 0.55 + i * 0.38
        txb(sl, "✓", px + 0.2, iy, 0.2, 0.3, size=9, color=TEAL)
        txb(sl, text, px + 0.45, iy, pw - 0.65, 0.3, size=9, color=WHITE)

    footer = spec.get("footer", "")
    if footer:
        txb(sl, footer, 0.5, 5.25, 9.0, 0.3, size=8, color=GRAY)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 14. Two-image slide — two screenshots side-by-side with captions
#     spec: {eyebrow?, title, left_image?, left_caption?,
#            right_image?, right_caption?, footer?}
# ─────────────────────────────────────────────────────────────────────────────
def render_two_image(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    eyebrow = spec.get("eyebrow", "")
    if eyebrow:
        set_ph(sl, 1, eyebrow, size=10, color=TEAL, italic=True)
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)

    for side, x in [("left", 0.5), ("right", 5.2)]:
        img = spec.get(f"{side}_image", "")
        cap = spec.get(f"{side}_caption", "")
        if img and os.path.isfile(img):
            sl.shapes.add_picture(img,
                                  Inches(x), Inches(1.6), Inches(4.3), Inches(3.4))
        if cap:
            txb(sl, cap, x, 4.95, 4.3, 0.3,
                size=8, color=GRAY, align=PP_ALIGN.CENTER)

    footer = spec.get("footer", "")
    if footer:
        txb(sl, footer, 0.5, 5.25, 9.0, 0.3, size=8, color=GRAY)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 15. Value propositions — stacked icon + bold title + description
#     spec: {eyebrow?, title, subtitle?, props: [{icon?, title, description}], footer?}
# ─────────────────────────────────────────────────────────────────────────────
def render_value_props(prs, SL, spec):
    sl = _new(prs, SL, "title_content")
    eyebrow = spec.get("eyebrow", "")
    if eyebrow:
        set_ph(sl, 1, eyebrow, size=10, color=TEAL, italic=True)
    set_ph(sl, 0, spec.get("title", ""),
           size=22, bold=True, color=WHITE)

    subtitle = spec.get("subtitle", "")
    if subtitle:
        txb(sl, subtitle, 0.6, 0.8, 9.0, 0.9,
            size=14, color=WHITE)

    props = spec.get("props", [])
    y_start = 1.95
    spacing = 0.6

    for i, p in enumerate(props):
        y = y_start + i * spacing
        icon = p.get("icon", "●")
        title = p.get("title", "")
        desc  = p.get("description", "")
        # Icon circle
        txb(sl, icon, 0.6, y, 0.3, 0.3, size=11, color=TEAL)
        # Bold title + description on same line
        txb(sl, title, 1.05, y, 2.5, 0.3, size=11, bold=True, color=WHITE)
        txb(sl, desc, 1.05, y + 0.25, 8.0, 0.3, size=9, color=GRAY)

    footer = spec.get("footer", "")
    if footer:
        txb(sl, footer, 0.5, 5.25, 9.0, 0.3, size=8, color=GRAY)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# 16. CTA / Call-to-action closing (brand, headline, sub-text, CTA button)
#     spec: {brand?, headline, sub_text?, cta_text?, footer?}
# ─────────────────────────────────────────────────────────────────────────────
def render_cta(prs, SL, spec):
    sl = _new(prs, SL, "title_center")
    brand = spec.get("brand", "dynatrace")
    txb(sl, brand, 0.6, 0.4, 3.0, 0.5, size=16, bold=False, color=WHITE)
    txb(sl, spec.get("headline", ""), 0.6, 2.0, 8.0, 0.9,
        size=32, bold=True, color=WHITE)
    sub = spec.get("sub_text", "")
    if sub:
        txb(sl, sub, 0.6, 3.0, 7.0, 0.8, size=14, color=GRAY)
    cta = spec.get("cta_text", "")
    if cta:
        btn = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.6), Inches(4.2), Inches(2.8), Inches(0.55))
        btn.fill.solid(); btn.fill.fore_color.rgb = TEAL
        btn.line.fill.background()
        tf = btn.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = cta
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# Dispatcher — given a slide spec dict, call the right renderer
# ─────────────────────────────────────────────────────────────────────────────
_RENDERERS = {
    "title":        render_title,
    "section":      render_section,
    "bullets":      render_bullets,
    "table":        render_table,
    "two_column":   render_two_column,
    "text":         render_text,
    "image":        render_image,
    "comparison":   render_comparison,
    "closing":      render_closing,
    "hero":         render_hero,
    "card_grid":    render_card_grid,
    "icon_bullets": render_icon_bullets,
    "split_panel":  render_split_panel,
    "two_image":    render_two_image,
    "value_props":  render_value_props,
    "cta":          render_cta,
}

def render_slide(prs, SL, spec: dict):
    """Render a single slide from a spec dict. Returns the slide object."""
    slide_type = spec.get("type", "bullets")
    renderer = _RENDERERS.get(slide_type)
    if renderer is None:
        raise ValueError(f"Unknown slide type: '{slide_type}'. "
                         f"Valid types: {list(_RENDERERS.keys())}")
    return renderer(prs, SL, spec)


def render_all(prs, SL, slides: list[dict]):
    """Render a list of slide specs, returning all slide objects."""
    results = []
    for spec in slides:
        results.append(render_slide(prs, SL, spec))
    return results
