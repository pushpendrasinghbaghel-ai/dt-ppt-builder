"""
Low-level drawing helpers for the DT PPT builder.
All positional args are in inches.
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

from .brand import (WHITE, TEAL, GREEN, ORANGE, GRAY, DGRAY, DDGRAY, DTDARK,
                    status_color, RGBColor)


# ─────────────────────────────────────────────────────────────────────────────
# Placeholder setter
# ─────────────────────────────────────────────────────────────────────────────
def set_ph(slide, ph_idx, text, size=None, bold=False, color=WHITE,
           align=PP_ALIGN.LEFT, italic=False):
    """Write text into a layout placeholder by idx."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = text
            if size:       r.font.size   = Pt(size)
            r.font.bold    = bold
            r.font.italic  = italic
            r.font.color.rgb = color
            return ph
    return None


# ─────────────────────────────────────────────────────────────────────────────
# Text box
# ─────────────────────────────────────────────────────────────────────────────
def txb(slide, text, l, t, w, h,
        size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    """Add a word-wrapped textbox."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return tb


# ─────────────────────────────────────────────────────────────────────────────
# Bulleted paragraph block
# ─────────────────────────────────────────────────────────────────────────────
def para_block(slide, lines, l, t, w, h, size=11, color=WHITE,
               hdr=None, hdr_color=TEAL, hdr_size=13):
    """Textbox with optional bold header then bullet lines (▸ prefix)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if hdr:
        p = tf.paragraphs[0]; first = False
        r = p.add_run()
        r.text           = hdr
        r.font.size      = Pt(hdr_size)
        r.font.bold      = True
        r.font.color.rgb = hdr_color
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before   = Pt(3)
        r = p.add_run()
        r.text           = f"\u25b8  {line}"
        r.font.size      = Pt(size)
        r.font.color.rgb = color
    return tb


# ─────────────────────────────────────────────────────────────────────────────
# Status badges bar  (✅ Now | ⚡ Partial | 🗺 Roadmap)
# ─────────────────────────────────────────────────────────────────────────────
def status_bar(slide, now, partial, roadmap, total, left=0.7, top=2.38):
    """Three coloured pill badges + total requirement count."""
    def _badge(text, l, bg):
        s = slide.shapes.add_shape(5, Inches(l), Inches(top), Inches(1.32), Inches(0.27))
        s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
        tf = s.text_frame
        tf.margin_left = tf.margin_right = Inches(0.04)
        tf.margin_top  = tf.margin_bottom = Inches(0.02)
        p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r  = p.add_run(); r.text = text
        r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = WHITE
    _badge(f"\u2705  {now} Now",         left,        GREEN)
    _badge(f"\u26a1  {partial} Partial",  left + 1.4,  ORANGE)
    _badge(f"\U0001f5fa  {roadmap} Roadmap", left + 2.8, RGBColor(0x55, 0x55, 0x55))
    txb(slide, f"of {total} requirements", left + 4.25, top + 0.02, 2.5, 0.28,
        size=10, color=GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# Requirements table  — 4 columns: Requirement | Description | Status | Signal
# ─────────────────────────────────────────────────────────────────────────────
def req_table(slide, reqs, l, t, w, h):
    """
    reqs: list of (Requirement, Description, Status, Signal)
    Status must contain ✅ / ⚡ / 🗺 for automatic colour coding.
    """
    cols = ["Requirement", "Description", "Status", "Signal"]
    tbl  = slide.shapes.add_table(
        len(reqs) + 1, 4, Inches(l), Inches(t), Inches(w), Inches(h)).table
    for c, cw in enumerate([w * f for f in [0.40, 0.29, 0.17, 0.14]]):
        tbl.columns[c].width = Inches(cw)
    # Header row
    for c, h_txt in enumerate(cols):
        cell = tbl.cell(0, c)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = h_txt
        r.font.size = Pt(8); r.font.bold = True; r.font.color.rgb = TEAL
        cell.fill.solid(); cell.fill.fore_color.rgb = DTDARK
    # Data rows
    for ri, (name, desc, st, sig) in enumerate(reqs):
        bg = DGRAY if ri % 2 == 0 else DDGRAY
        for c, val in enumerate([name, desc, st, sig]):
            cell = tbl.cell(ri + 1, c)
            p    = cell.text_frame.paragraphs[0]
            r    = p.add_run(); r.text = val
            r.font.size = Pt(7.5)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            r.font.color.rgb = (status_color(val) if c == 2 else
                                TEAL               if c == 3 else WHITE)
    return tbl


# ─────────────────────────────────────────────────────────────────────────────
# Image helper
# ─────────────────────────────────────────────────────────────────────────────
def add_img(slide, path, l, t, w, h):
    """Embed an image if the file exists; silently skip if not."""
    if path and os.path.exists(path):
        return slide.shapes.add_picture(path, Inches(l), Inches(t),
                                        Inches(w), Inches(h))
    print(f"  [WARN] image not found: {path}")
    return None


# ─────────────────────────────────────────────────────────────────────────────
# Coverage matrix table  (Summary slide)
# ─────────────────────────────────────────────────────────────────────────────
def coverage_table(slide, domains, l=0.7, t=2.3, w=11.94, h=4.6):
    """
    domains: list of dicts with keys:
        name, total, now, partial, roadmap
    Final row is auto-generated TOTAL.
    """
    tot_total   = sum(d["total"]   for d in domains)
    tot_now     = sum(d["now"]     for d in domains)
    tot_partial = sum(d["partial"] for d in domains)
    tot_roadmap = sum(d["roadmap"] for d in domains)

    data_rows = [[d["name"], str(d["total"]),
                  str(d["now"]), str(d["partial"]), str(d["roadmap"])]
                 for d in domains]
    data_rows.append([
        "TOTAL", str(tot_total),
        f"{tot_now} ({round(tot_now/tot_total*100)}%)",
        f"{tot_partial} ({round(tot_partial/tot_total*100)}%)",
        f"{tot_roadmap} ({round(tot_roadmap/tot_total*100)}%)",
    ])

    tbl = slide.shapes.add_table(
        len(data_rows) + 1, 5, Inches(l), Inches(t), Inches(w), Inches(h)).table
    for c, cw in enumerate([5.5, 1.1, 1.8, 1.8, 1.74]):
        tbl.columns[c].width = Inches(cw)

    from .brand import GREEN, ORANGE, GRAY, WHITE, DTDARK, DGRAY, DDGRAY, TEAL
    hdrs = ["Domain", "Total", "\u2705 Now", "\u26a1 Partial", "\U0001f5fa Roadmap"]
    for c, h_txt in enumerate(hdrs):
        cell = tbl.cell(0, c)
        p    = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r    = p.add_run(); r.text = h_txt
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = TEAL
        cell.fill.solid(); cell.fill.fore_color.rgb = DTDARK

    for ri, row in enumerate(data_rows):
        is_total = (ri == len(data_rows) - 1)
        bg = DTDARK if is_total else (DGRAY if ri % 2 == 0 else DDGRAY)
        for c, val in enumerate(row):
            cell = tbl.cell(ri + 1, c)
            p    = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r    = p.add_run(); r.text = val
            r.font.size = Pt(11 if is_total else 10)
            r.font.bold = is_total
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            if c == 2:   r.font.color.rgb = GREEN
            elif c == 3: r.font.color.rgb = ORANGE
            elif c == 4: r.font.color.rgb = GRAY
            else:        r.font.color.rgb = WHITE
    return tbl
